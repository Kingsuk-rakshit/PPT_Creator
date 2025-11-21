import os
import json
import requests
import math
import time
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from dotenv import load_dotenv
import uuid


# Load environment variables from .env file immediately
load_dotenv()


# Global defaults from Environment Variables
client = None
# Default to a common Groq model, user can override in .env
model_name = os.getenv("GROQ_MODEL_NAME")
pexels_api_key = os.getenv("PEXELS_API_KEY")


def initialize_client():
    """
    Attempts to initialize the Groq client using Environment Variables.
    """
    global client
    api_key = os.getenv("GROQ_API_KEY")
   
    if api_key:
        try:
            client = Groq(api_key=api_key)
            return True
        except Exception as e:
            print(f"Groq Init Failed: {e}")
            return False
    return False


# Try to initialize immediately
initialize_client()


# AGENT 1 - CONTENT GENERATOR
def generate_slide_content(topic, feedback=None, current_content=None):
    if not client:
        # Try one last time to init from env if client is missing
        if not initialize_client():
            return "Error: Groq Client not initialized. Check your .env file for GROQ_API_KEY."


    system_prompt = """
    You are a Professional PowerPoint Content Generator.
    Your goal is to create a detailed and informative presentation.
   
    Structure Requirements:
    1.  **Title Slide**: A catchy title for the presentation.
    2.  **Table of Contents**: A list of the titles of ALL content slides.
    3.  **Content Slides**: Create 5 content slides by default. If the user requests a specific number of slides, generate that many. Each slide must have a title and exactly 3 descriptive bullet points.
    4.  **Conclusion Slide**: A summary of the key takeaways. Provide exactly 3 concise bullet points.
   
    Output strictly in JSON format with this structure:
    {
        "presentation_title": "Main Title of Presentation",
        "table_of_contents": ["Slide 1 Title", "Slide 2 Title", "Slide 3 Title", "Slide 4 Title", "Slide 5 Title"],
        "slides": [
            {
                "title": "Slide Title",
                "content": [
                    "Detailed bullet point 1.",
                    "Detailed bullet point 2.",
                    "Detailed bullet point 3."
                ],
                "image_description": "A search query for a stock photo website (e.g., 'business meeting', 'nature landscape')."
            }
        ],
        "conclusion": {
            "title": "Conclusion",
            "content": ["Concise takeaway 1", "Concise takeaway 2", "Concise takeaway 3"]
        }
    }
    Do not include markdown formatting like ```json.
    """
   
    user_prompt = f"Create a presentation about: {topic}"
   
    if feedback and current_content:
        user_prompt = f"""
        Current Content: {current_content}
        User Feedback: {feedback}
        Refine the content based on the feedback. Return the full updated JSON.
        """


    try:
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            response_format={"type": "json_object"},
            temperature=0.7
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {str(e)}"


# AGENT 2 - IMAGE GENERATOR
def generate_image(prompt):
    """Uses Pexels API to find an image URL"""
    if not pexels_api_key:
        print("Error: PEXELS_API_KEY not found in environment variables.")
        return None


    print(f"Searching Pexels for: '{prompt}'...")
    try:
        headers = {
            'Authorization': pexels_api_key,
            'User-Agent': 'PPTCreator/1.0'
        }
        params = {
            'query': prompt,
            'per_page': 1,
            'orientation': 'landscape',
            'size': 'large'
        }
        response = requests.get('https://api.pexels.com/v1/search', headers=headers, params=params, timeout=10)
       
        if response.status_code == 200:
            data = response.json()
            if data.get('photos'):
                image_url = data['photos'][0]['src']['landscape']
                print(f"Image found: {image_url}")
                return image_url
            else:
                print(f"No images found on Pexels for: '{prompt}'")
                return None
        else:
            print(f"Pexels API Error: {response.status_code} - {response.text}")
            return None
    except Exception as e:
        print(f"Error searching Pexels: {e}")
        return None


# PPT CREATOR - COMBINING THE EXTRACTED RESOURCES INTO ONE PRESENTATION
def create_ppt_file(slide_data, include_images=True, theme_color='#003366'):
    try:
        data = json.loads(slide_data)
    except json.JSONDecodeError:
        return None


    prs = Presentation()
    # Set slide dimensions for widescreen 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)


    # Helper to convert hex to RGB
    def hex_to_rgb(hex_color):
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


    # Theme Colors
    try:
        r, g, b = hex_to_rgb(theme_color)
        NAVY_BLUE = RGBColor(r, g, b) # User selected color
    except:
        NAVY_BLUE = RGBColor(0, 51, 102) # Fallback


    DARK_GRAY = RGBColor(80, 80, 80)


    def fetch_image(url, retries=3, timeout=30):
        print(f"Downloading image from: {url}")
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}
        for attempt in range(retries):
            try:
                resp = requests.get(url, timeout=timeout, headers=headers)
                if resp.status_code == 200:
                    print("Image downloaded successfully.")
                    return resp.content
                else:
                    print(f"Image download failed: {resp.status_code}")
            except Exception as e:
                print(f"Image download error (Attempt {attempt+1}): {e}")
            time.sleep(1 * (attempt + 1))
        print("Failed to download image after retries.")
        return None


    # 1. Title Slide
    slide_layout = prs.slide_layouts[6] # Blank for custom
    slide = prs.slides.add_slide(slide_layout)
   
    # Decorative Bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_BLUE
    shape.line.fill.background()


    # Title
    title_text = data.get('presentation_title', 'Presentation')
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title_text
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.alignment = PP_ALIGN.CENTER
   
    # Dynamic Title Sizing
    if len(title_text) > 40:
        p.font.size = Pt(32)
    elif len(title_text) > 25:
        p.font.size = Pt(38)
    else:
        p.font.size = Pt(44)


    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    tf = sub_box.text_frame
    p = tf.add_paragraph()
    p.text = "Generated by AI Agent"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER


    # 2. Table of Contents
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
   
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Table of Contents"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = NAVY_BLUE
    p.alignment = PP_ALIGN.CENTER
   
    # Separator Line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.5), Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY_BLUE
    line.line.fill.background()


    toc_items = data.get('table_of_contents', [])
    item_count = len(toc_items)
   
    # Dynamic ToC Sizing
    if item_count > 8:
        toc_font_size = Pt(14)
        toc_spacing = Pt(8)
    elif item_count > 5:
        toc_font_size = Pt(16)
        toc_spacing = Pt(12)
    else:
        toc_font_size = Pt(20)
        toc_spacing = Pt(14)


    # Aggressive 2-column switch to prevent overflow
    if item_count > 4:
        # Two Column Layout
        mid = math.ceil(item_count / 2)
        col1_items = toc_items[:mid]
        col2_items = toc_items[mid:]
       
        # Column 1
        box1 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(3))
        tf1 = box1.text_frame
        tf1.word_wrap = True
        for item in col1_items:
            p = tf1.add_paragraph()
            p.text = f"• {item}"
            p.font.size = toc_font_size
            p.space_after = toc_spacing
            p.font.color.rgb = DARK_GRAY
           
        # Column 2
        box2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(3))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        for item in col2_items:
            p = tf2.add_paragraph()
            p.text = f"• {item}"
            p.font.size = toc_font_size
            p.space_after = toc_spacing
            p.font.color.rgb = DARK_GRAY
    else:
        # Single Column Centered
        box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(3))
        tf = box.text_frame
        tf.word_wrap = True
        for item in toc_items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = toc_font_size
            p.space_after = toc_spacing
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.LEFT


    # 3. Content Slides
    slides_content = data.get('slides', [])
    for i, slide_info in enumerate(slides_content):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank


        # Title
        slide_title = slide_info.get('title', 'Slide')
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tframe = title_shape.text_frame
        tframe.word_wrap = True
        tp = tframe.paragraphs[0]
        tp.text = slide_title
        tp.font.bold = True
        tp.font.color.rgb = NAVY_BLUE
       
        # Dynamic Title Font
        if len(slide_title) > 50:
            tp.font.size = Pt(24)
        elif len(slide_title) > 35:
            tp.font.size = Pt(26)
        else:
            tp.font.size = Pt(28)
       
        # Separator Line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = NAVY_BLUE
        line.line.fill.background()


        # Try image download first
        img_data = None
        img_url = None
        if include_images and slide_info.get("image_description"):
            img_url = generate_image(slide_info['image_description'])
            if img_url:
                img_data = fetch_image(img_url)


        has_image = img_data is not None
       
        # Content Layout
        if has_image:
            content_width = 5.0
            image_x = Inches(6.0)
            image_y = Inches(2.0) # Moved down
            image_w = Inches(3.5)
        else:
            content_width = 9.0 # Full width if no image


        # Content box
        content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(content_width), Inches(3.9))
        tf = content_shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP


        bullets = [b.strip() for b in slide_info.get('content', []) if b.strip()]
        if not bullets:
            bullets = ["(No content generated)"]


        # Calculate total text length to determine font size
        total_chars = sum(len(b) for b in bullets)
       
        if total_chars > 600:
            content_font_size = Pt(12)
            content_spacing = Pt(4)
        elif total_chars > 400:
            content_font_size = Pt(14)
            content_spacing = Pt(6)
        else:
            content_font_size = Pt(16)
            content_spacing = Pt(8)


        for idx, bullet in enumerate(bullets):
            p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.level = 0
            p.font.size = content_font_size
            p.font.color.rgb = DARK_GRAY
            p.space_before = content_spacing
            p.space_after = content_spacing


        # Image placement
        if has_image:
            try:
                temp_filename = f"temp_img_{i}.jpg"
                with open(temp_filename, 'wb') as f:
                    f.write(img_data)
                # Add a border to image
                pic = slide.shapes.add_picture(temp_filename, image_x, image_y, width=image_w)
                line = pic.line
                line.color.rgb = NAVY_BLUE
                line.width = Pt(2)
                os.remove(temp_filename)
            except Exception as e:
                print(f"Image placement failed: {e}")


    # 4. Conclusion Slide
    conclusion_data = data.get('conclusion', {})
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
   
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tframe = title_shape.text_frame
    tp = tframe.paragraphs[0]
    tp.text = conclusion_data.get('title', 'Conclusion')
    tp.font.bold = True
    tp.font.size = Pt(28)
    tp.font.color.rgb = NAVY_BLUE
   
    # Separator Line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY_BLUE
    line.line.fill.background()
   
    # Content
    content_shape = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3.5))
    tf = content_shape.text_frame
    tf.word_wrap = True
    conc_points = conclusion_data.get('content', [])
   
    # Dynamic Conclusion Sizing
    total_conc_chars = sum(len(p) for p in conc_points)
    if total_conc_chars > 500:
        conc_font_size = Pt(14)
        conc_spacing = Pt(8)
    elif total_conc_chars > 300:
        conc_font_size = Pt(16)
        conc_spacing = Pt(10)
    else:
        conc_font_size = Pt(20)
        conc_spacing = Pt(14)


    for point in conc_points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.level = 0
        p.font.size = conc_font_size
        p.font.color.rgb = DARK_GRAY
        p.space_after = conc_spacing


    # 5. Thank You Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
   
    # Full background color for Thank You slide
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY_BLUE
   
    # Center text box
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Thank You"
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = RGBColor(255, 255, 255) # White text
    p.alignment = PP_ALIGN.CENTER


    unique_id = uuid.uuid4().hex[:6]
    filename = f"generated_presentation_{unique_id}.pptx"
    prs.save(filename)
    return filename
